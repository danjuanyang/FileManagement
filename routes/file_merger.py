# file_merger.py
import os
import re
import shutil
import tempfile
from datetime import datetime
import io

from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.opc.oxml import qn
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfMerger
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase.pdfmetrics import registerFontFamily

from flask import current_app
from sqlalchemy import or_, text, func

# 导入需要的模型
from models import db, Project, ProjectFile, ProjectStage, User, StageTask


def setup_fonts():
    """设置字体"""
    try:
        # 获取字体文件路径
        font_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'fonts', 'simsun.ttf')

        # 注册基本字体
        pdfmetrics.registerFont(TTFont('SimSun', font_path))

        # 注册字体变体
        pdfmetrics.registerFontFamily(
            'SimSun',
            normal='SimSun',
            bold='SimSun',  # 使用同一个字体文件
            italic='SimSun',  # 使用同一个字体文件
            boldItalic='SimSun'  # 使用同一个字体文件
        )
        return True
    except Exception as e:
        print(f"字体设置失败: {str(e)}")
        return False


def create_pdf_style():
    """创建PDF样式"""
    styles = getSampleStyleSheet()

    # 确保使用已注册的字体
    default_font = 'SimSun'

    # 基础样式
    basic_style = ParagraphStyle(
        'BasicStyle',
        parent=styles['Normal'],
        fontName=default_font,
        fontSize=12,
        leading=14,
        allowWidows=0,
        allowOrphans=0
    )

    # 标题样式
    title_style = ParagraphStyle(
        'TitleStyle',
        parent=basic_style,
        fontName=default_font,
        fontSize=16,
        leading=20,
        spaceAfter=30,
        alignment=1
    )

    # 居中样式
    center_style = ParagraphStyle(
        'CenterStyle',
        parent=basic_style,
        fontName=default_font,
        alignment=1
    )

    # 右对齐样式
    right_style = ParagraphStyle(
        'RightStyle',
        parent=basic_style,
        fontName=default_font,
        alignment=2
    )

    return {
        'basic': basic_style,
        'title': title_style,
        'center': center_style,
        'right': right_style
    }


def extract_prefix_number(filename):
    """从文件名中提取前缀数字"""
    match = re.match(r'^(\d+)', filename)
    if match:
        return int(match.group(1))
    return float('inf')


def sort_files_by_prefix(files):
    """按照文件名前缀数字排序文件"""
    return sorted(files, key=lambda x: extract_prefix_number(x.file_name))


def convert_word_to_pdf(input_file, output_file):
    """将Word文档转换为PDF，支持横向格式"""
    try:
        if not setup_fonts():
            raise Exception("字体设置失败")

        doc = Document(input_file)

        # 检查文档方向
        section = doc.sections[0]
        is_landscape = section.orientation == 1  # 1 表示横向

        # 根据方向设置页面大小
        pagesize = A4 if not is_landscape else landscape(A4)

        pdf = SimpleDocTemplate(
            output_file,
            pagesize=pagesize,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )

        styles = create_pdf_style()
        story = []

        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            # 根据段落样式选择对应的PDF样式
            if paragraph.style.name.startswith('Heading'):
                p = Paragraph(text, styles['title'])
            else:
                # 根据对齐方式选择样式
                if paragraph.alignment == WD_PARAGRAPH_ALIGNMENT.CENTER:
                    p = Paragraph(text, styles['center'])
                elif paragraph.alignment == WD_PARAGRAPH_ALIGNMENT.RIGHT:
                    p = Paragraph(text, styles['right'])
                else:
                    p = Paragraph(text, styles['basic'])

            story.append(p)
            story.append(Spacer(1, 12))

        pdf.build(story)
        return True
    except Exception as e:
        print(f"转换Word文档失败: {str(e)}")
        return False


def should_use_landscape(data, max_content_width):
    """基于内容宽度判断是否使用横向布局"""
    if not data:
        return False

    # 计算预估的内容总宽度（考虑到字体大小和页面边距）
    estimated_width = max_content_width * 7  # 假设每个字符平均7个单位宽度

    # A4纸的宽度（portrait模式）减去边距
    A4_PORTRAIT_WIDTH = A4[0] - 30  # 减去左右各15的边距

    # 如果预估宽度超过可用宽度的75%，建议使用横向
    return estimated_width > (A4_PORTRAIT_WIDTH * 0.75)


def format_number(value, number_format):
    """格式化数字，处理常见的Excel数字格式"""
    try:
        # 处理百分比格式
        if '%' in number_format:
            return f"{value:.2%}"
        # 处理金额格式
        elif '¥' in number_format or '$' in number_format:
            return f"{value:,.2f}"
        # 处理小数位数
        elif '0' in number_format:
            decimal_places = len(number_format.split('.')[-1]) if '.' in number_format else 0
            return f"{value:,.{decimal_places}f}"
        # 默认格式
        return str(value)
    except:
        return str(value)


def add_page_number(canvas, doc):
    """添加页码"""
    try:
        canvas.saveState()
        canvas.setFont('SimSun', 8)
        page_num = canvas.getPageNumber()
        text = f"第 {page_num} 页"
        canvas.drawRightString(doc.pagesize[0] - 20, 20, text)
        canvas.restoreState()
    except Exception as e:
        print(f"添加页码失败: {str(e)}")
        return False


def convert_excel_to_pdf(input_file, output_file):
    """将Excel文档转换为PDF，优化公式显示和宽表格处理"""
    try:
        if not setup_fonts():
            raise Exception("字体设置失败")

        wb = load_workbook(input_file, data_only=True)  # 添加data_only=True以获取公式计算结果

        # 遍历所有工作表，找出是否有任何一个工作表需要横向布局
        needs_landscape = False
        sheet_orientations = {}  # 存储每个工作表的方向

        for sheet in wb.worksheets:
            # 收集工作表数据，同时计算实际内容宽度
            sheet_data = []
            max_content_width = 0

            for row in sheet.rows:
                row_data = []
                for cell in row:
                    # 获取单元格的实际值而不是公式
                    if cell.value is not None:
                        value = cell.value
                        # 处理日期类型
                        if isinstance(value, datetime):
                            value = value.strftime('%Y-%m-%d %H:%M:%S')
                        # 处理数字格式
                        elif isinstance(value, (int, float)):
                            if cell.number_format != 'General':
                                try:
                                    value = format_number(value, cell.number_format)
                                except:
                                    value = str(value)
                            else:
                                value = str(value)
                        else:
                            value = str(value)
                    else:
                        value = ''

                    row_data.append(value)
                    # 计算内容宽度（考虑中文字符）
                    content_width = sum(2 if '\u4e00' <= char <= '\u9fff' else 1 for char in str(value))
                    max_content_width = max(max_content_width, content_width)

                sheet_data.append(row_data)

            # 基于内容宽度判断方向
            sheet_needs_landscape = should_use_landscape(sheet_data, max_content_width)
            sheet_orientations[sheet.title] = sheet_needs_landscape
            needs_landscape = needs_landscape or sheet_needs_landscape

        # 设置PDF尺寸
        pagesize = landscape(A4) if needs_landscape else A4

        pdf = SimpleDocTemplate(
            output_file,
            pagesize=pagesize,
            rightMargin=15,  # 减小边距以获得更多空间
            leftMargin=15,
            topMargin=15,
            bottomMargin=15
        )

        elements = []
        styles = create_pdf_style()

        for sheet in wb.worksheets:
            # 添加分页符（除了第一个工作表）
            if elements:
                elements.append(PageBreak())

            # 添加工作表标题
            elements.append(Paragraph(sheet.title, styles['title']))
            elements.append(Spacer(1, 10))

            # 获取数据并计算最佳列宽
            data = []
            max_col_widths = []

            for row in sheet.rows:
                row_data = []
                for cell in row:
                    # 获取单元格的实际值
                    if cell.value is not None:
                        value = cell.value
                        # 处理日期类型
                        if isinstance(value, datetime):
                            value = value.strftime('%Y-%m-%d %H:%M:%S')
                        # 处理数字格式
                        elif isinstance(value, (int, float)):
                            if cell.number_format != 'General':
                                try:
                                    value = format_number(value, cell.number_format)
                                except:
                                    value = str(value)
                            else:
                                value = str(value)
                        else:
                            value = str(value)
                    else:
                        value = ''

                    row_data.append(value)

                data.append(row_data)

                # 更新列宽度
                while len(max_col_widths) < len(row_data):
                    max_col_widths.append(0)
                for i, value in enumerate(row_data):
                    # 计算内容宽度（考虑中文字符）
                    content_width = sum(2 if '\u4e00' <= char <= '\u9fff' else 1 for char in str(value))
                    max_col_widths[i] = max(max_col_widths[i], content_width)

            if data:
                # 计算可用宽度和列宽
                available_width = pagesize[0] - pdf.leftMargin - pdf.rightMargin
                total_units = sum(max_col_widths)

                # 计算最小列宽和理想列宽
                min_col_width = 20  # 最小列宽
                col_widths = []

                for width in max_col_widths:
                    if total_units > 0:
                        # 按比例分配宽度，但确保最小宽度
                        col_width = max((width / total_units) * available_width, min_col_width)
                    else:
                        col_width = available_width / len(max_col_widths)
                    col_widths.append(col_width)

                # 创建表格
                table = Table(
                    data,
                    colWidths=col_widths,
                    style=[
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('FONTNAME', (0, 0), (-1, -1), 'SimSun'),
                        ('FONTSIZE', (0, 0), (-1, -1), 7),  # 减小字体以适应更多内容
                        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                        ('WORDWRAP', (0, 0), (-1, -1), True),
                        ('LEFTPADDING', (0, 0), (-1, -1), 2),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 2),
                        ('TOPPADDING', (0, 0), (-1, -1), 2),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
                    ]
                )
                elements.append(table)

        # 构建PDF
        pdf.build(elements, onFirstPage=add_page_number, onLaterPages=add_page_number)
        return True

    except Exception as e:
        print(f"转换Excel文档失败: {str(e)}")
        return False


def convert_ppt_to_pdf(input_file, output_file):
    """将PPT文档转换为PDF，支持横向格式"""
    try:
        prs = Presentation(input_file)

        # PowerPoint 默认是横向的，所以使用横向A4
        pagesize = landscape(A4)
        width, height = pagesize

        c = canvas.Canvas(output_file, pagesize=pagesize)

        for slide in prs.slides:
            y = height - 50  # 起始位置
            x = 50  # 左边距

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # 分行处理文本
                    for line in shape.text.split('\n'):
                        if line.strip():
                            # 如果文本太长需要换行
                            words = line.split()
                            current_line = []
                            for word in words:
                                test_line = ' '.join(current_line + [word])
                                if c.stringWidth(test_line) < (width - 100):  # 留出左右边距
                                    current_line.append(word)
                                else:
                                    if current_line:
                                        c.drawString(x, y, ' '.join(current_line))
                                        y -= 20
                                        current_line = [word]
                                    else:
                                        # 如果单词太长，强制换行
                                        c.drawString(x, y, word)
                                        y -= 20

                            if current_line:
                                c.drawString(x, y, ' '.join(current_line))
                                y -= 20

                if y < 50:  # 如果页面空间不足，创建新页面
                    c.showPage()
                    y = height - 50

            c.showPage()  # 每个幻灯片结束后创建新页面

        c.save()
        return True
    except Exception as e:
        print(f"转换PPT文档失败: {str(e)}")
        return False


def convert_to_pdf(input_file, output_file):
    """将Office文件转换为PDF"""
    ext = os.path.splitext(input_file)[1].lower()

    try:
        if ext in ['.doc', '.docx']:
            return convert_word_to_pdf(input_file, output_file)
        elif ext in ['.xls', '.xlsx']:
            return convert_excel_to_pdf(input_file, output_file)
        elif ext in ['.ppt', '.pptx']:
            return convert_ppt_to_pdf(input_file, output_file)
        return False
    except Exception as e:
        print(f"转换文件失败 {input_file}: {str(e)}")
        return False


def append_word_content(target_doc, word_path):
    """将Word文档内容追加到目标文档"""
    try:
        # 打开源文档
        source_doc = Document(word_path)

        # 复制段落
        for paragraph in source_doc.paragraphs:
            if paragraph.text.strip():  # 跳过空段落
                # 复制文本和格式
                p = target_doc.add_paragraph()
                for run in paragraph.runs:
                    new_run = p.add_run(run.text)
                    # 复制格式
                    new_run.bold = run.bold
                    new_run.italic = run.italic
                    new_run.underline = run.underline
                    if run.font.color.rgb:
                        new_run.font.color.rgb = run.font.color.rgb

        # 复制表格
        for table in source_doc.tables:
            # 创建新表格
            rows = len(table.rows)
            cols = len(table.columns)
            new_table = target_doc.add_table(rows=rows, cols=cols)
            new_table.style = table.style

            # 复制单元格内容
            for i, row in enumerate(table.rows):
                for j, cell in enumerate(row.cells):
                    new_cell = new_table.cell(i, j)
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():  # 跳过空段落
                            p = new_cell.paragraphs[0] if j == 0 and i == 0 else new_cell.add_paragraph()
                            for run in paragraph.runs:
                                new_run = p.add_run(run.text)
                                new_run.bold = run.bold
                                new_run.italic = run.italic

        # 添加空行
        target_doc.add_paragraph()

    except Exception as e:
        # 添加错误提示
        target_doc.add_paragraph(f"无法完全复制Word文档内容: {str(e)}")


def append_excel_summary(target_doc, excel_path):
    """添加Excel内容摘要到目标文档"""
    try:
        # 加载Excel文件
        wb = load_workbook(excel_path, data_only=True)

        for sheet in wb.worksheets:
            # 添加工作表名称
            target_doc.add_heading(f'工作表: {sheet.title}', level=3)

            # 获取数据范围
            max_row = min(sheet.max_row, 100)  # 限制最多显示100行
            max_col = min(sheet.max_column, 15)  # 限制最多显示15列

            # 创建表格
            if max_row > 0 and max_col > 0:
                table = target_doc.add_table(rows=max_row, cols=max_col)
                table.style = 'Table Grid'

                # 填充数据
                for i in range(max_row):
                    for j in range(max_col):
                        cell_value = sheet.cell(row=i + 1, column=j + 1).value
                        if cell_value is not None:
                            # 处理特殊类型
                            if isinstance(cell_value, datetime):
                                cell_value = cell_value.strftime('%Y-%m-%d %H:%M:%S')
                            else:
                                cell_value = str(cell_value)

                            # 限制单元格内容长度
                            if len(cell_value) > 100:
                                cell_value = cell_value[:100] + "..."

                            table.cell(i, j).text = cell_value

                # 标记表头
                for j in range(max_col):
                    for run in table.cell(0, j).paragraphs[0].runs:
                        run.bold = True
            else:
                target_doc.add_paragraph("(工作表为空)")

            # 添加说明
            if sheet.max_row > max_row or sheet.max_column > max_col:
                target_doc.add_paragraph(
                    f"注: 完整工作表包含 {sheet.max_row} 行, {sheet.max_column} 列，此处仅显示部分内容。")

            # 添加分隔符
            target_doc.add_paragraph()

    except Exception as e:
        target_doc.add_paragraph(f"无法读取Excel内容: {str(e)}")


def append_ppt_summary(target_doc, ppt_path):
    """添加PPT内容摘要到目标文档"""
    try:
        # 加载PPT文件
        prs = Presentation(ppt_path)

        # 添加摘要信息
        target_doc.add_paragraph(f"PowerPoint演示文稿: 共 {len(prs.slides)} 张幻灯片")

        # 遍历幻灯片
        for i, slide in enumerate(prs.slides):
            # 添加幻灯片标题
            slide_title = "无标题"
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip() and hasattr(shape, 'is_title') and shape.is_title:
                    slide_title = shape.text.strip()
                    break

            title = target_doc.add_heading(f"幻灯片 {i + 1}: {slide_title}", level=3)

            # 提取幻灯片内容
            content = target_doc.add_paragraph()
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text and text != slide_title:  # 避免重复标题
                        content.add_run(text + "\n\n")

            # 如果没有提取到内容
            if not content.text.strip():
                content.add_run("(此幻灯片无文本内容或仅包含标题)")

            # 添加分隔线
            if i < len(prs.slides) - 1:
                target_doc.add_paragraph("---")

    except Exception as e:
        target_doc.add_paragraph(f"无法读取PowerPoint内容: {str(e)}")


def merge_project_files_to_pdf(project_id):
    """将项目文件按任务顺序合并为PDF"""
    temp_dir = None
    try:
        # 获取项目信息
        project = Project.query.get(project_id)
        if not project:
            return None, "项目不存在"

        # 创建临时目录
        temp_dir = tempfile.mkdtemp()
        os.makedirs(temp_dir, exist_ok=True)

        final_pdf_merger = PdfMerger()
        has_files = False
        processed_files = []
        skipped_files = []

        # 处理项目文件
        for stage in project.stages:
            for task in stage.tasks:
                files = ProjectFile.query.filter_by(
                    project_id=project_id,
                    stage_id=stage.id,
                    task_id=task.id
                ).all()

                sorted_files = sort_files_by_prefix(files)
                task_pdf_merger = PdfMerger()
                task_has_files = False

                for file in sorted_files:
                    ext = os.path.splitext(file.file_name)[1].lower()
                    supported_types = ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf']

                    if ext not in supported_types:
                        skipped_files.append({
                            'name': file.original_name,
                            'reason': '不支持的文件类型'
                        })
                        continue

                    input_path = os.path.join(current_app.root_path, file.file_path)
                    if not os.path.exists(input_path):
                        skipped_files.append({
                            'name': file.original_name,
                            'reason': '文件不存在'
                        })
                        continue

                    try:
                        temp_pdf = os.path.join(temp_dir, f"{file.id}.pdf")

                        if ext == '.pdf':
                            shutil.copy2(input_path, temp_pdf)
                            success = True
                        else:
                            success = convert_to_pdf(input_path, temp_pdf)

                        if success:
                            task_pdf_merger.append(temp_pdf)
                            task_has_files = True
                            processed_files.append(file.original_name)
                        else:
                            skipped_files.append({
                                'name': file.original_name,
                                'reason': '转换失败'
                            })
                    except Exception as e:
                        skipped_files.append({
                            'name': file.original_name,
                            'reason': f'处理错误: {str(e)}'
                        })

                if task_has_files:
                    task_pdf_path = os.path.join(temp_dir, f"task_{task.id}_merged.pdf")
                    task_pdf_merger.write(task_pdf_path)
                    task_pdf_merger.close()
                    final_pdf_merger.append(task_pdf_path)
                    has_files = True

        if not has_files:
            return None, "没有可合并的文件"

        # 生成最终PDF
        output_filename = f"{project.name}_merged.pdf"
        output_path = os.path.join(temp_dir, output_filename)
        final_pdf_merger.write(output_path)
        final_pdf_merger.close()

        # 复制到临时文件并返回
        temp_output = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        shutil.copy2(output_path, temp_output.name)

        return temp_output.name, None

    except Exception as e:
        return None, str(e)
    finally:
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                print(f"清理临时文件失败: {str(e)}")


def merge_project_files_to_word(project_id):
    """将项目文件按阶段和任务顺序合并为Word文档，每个阶段后插入一个空白页"""
    temp_dir = None
    try:
        print(f"开始合并项目 ID:{project_id} 的文件到Word...")

        # 获取项目信息
        project = Project.query.get(project_id)
        if not project:
            print(f"项目不存在: {project_id}")
            return None, "项目不存在"

        # 创建临时目录
        temp_dir = tempfile.mkdtemp()
        os.makedirs(temp_dir, exist_ok=True)
        print(f"临时目录已创建: {temp_dir}")

        # 创建主Word文档
        merged_doc = Document()

        # 设置中文字体
        style = merged_doc.styles['Normal']
        style.font.name = '微软雅黑'
        # 设置字体

        # 设置标题字体
        for i in range(1, 4):
            heading_style = merged_doc.styles[f'Heading {i}']
            heading_style.font.name = '微软雅黑'

        # 添加文档标题
        title = merged_doc.add_heading(f'{project.name} - 项目文档合集', 0)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        processed_files = []
        skipped_files = []

        print(f"开始处理项目 {project.name} 的各阶段文件...")
        # 处理各阶段和任务的文件
        for stage in project.stages:
            print(f"处理阶段: {stage.name}")
            # 添加阶段标题
            stage_title = merged_doc.add_heading(f'阶段：{stage.name}', level=1)
            stage_title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

            for task in stage.tasks:
                print(f"处理任务: {task.name}")
                # 添加任务标题
                task_title = merged_doc.add_heading(f'任务：{task.name}', level=2)
                task_title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

                # 获取该任务下的所有文件
                files = ProjectFile.query.filter_by(
                    project_id=project_id,
                    stage_id=stage.id,
                    task_id=task.id
                ).all()
                print(f"找到文件数量: {len(files)}")

                # 按文件名前缀排序
                sorted_files = sort_files_by_prefix(files)

                if not sorted_files:
                    merged_doc.add_paragraph('此任务下无文件')
                    continue

                # 处理每个文件
                for file in sorted_files:
                    try:
                        print(f"处理文件: {file.original_name}")
                        # 获取文件路径和扩展名
                        input_path = os.path.join(current_app.root_path, file.file_path)
                        ext = os.path.splitext(file.file_name)[1].lower()

                        if not os.path.exists(input_path):
                            error_msg = '文件不存在'
                            print(f"跳过文件 {file.original_name}: {error_msg}")
                            skipped_files.append({
                                'name': file.original_name,
                                'reason': error_msg
                            })
                            continue

                        # 添加文件信息
                        file_info = merged_doc.add_paragraph()
                        file_info.add_run(f'文件名: {file.original_name}').bold = True
                        file_info.add_run(f'\n上传时间: {file.upload_date.strftime("%Y-%m-%d %H:%M:%S")}')

                        # 根据文件类型处理
                        if ext in ['.doc', '.docx']:
                            print(f"合并Word文档内容: {file.original_name}")
                            # 直接合并Word文档内容
                            append_word_content(merged_doc, input_path)
                            processed_files.append(file.original_name)

                        elif ext in ['.xlsx', '.xls']:
                            print(f"添加Excel内容摘要: {file.original_name}")
                            # 添加Excel内容摘要
                            append_excel_summary(merged_doc, input_path)
                            processed_files.append(file.original_name)

                        elif ext in ['.pptx', '.ppt']:
                            print(f"添加PPT内容摘要: {file.original_name}")
                            # 添加PPT内容摘要
                            append_ppt_summary(merged_doc, input_path)
                            processed_files.append(file.original_name)

                        elif ext == '.pdf':
                            print(f"添加PDF引用: {file.original_name}")
                            # 添加PDF引用
                            merged_doc.add_paragraph(f'PDF文件: {file.original_name} (需单独查看)')
                            processed_files.append(file.original_name)

                        else:
                            error_msg = '不支持的文件类型'
                            print(f"跳过文件 {file.original_name}: {error_msg}")
                            skipped_files.append({
                                'name': file.original_name,
                                'reason': error_msg
                            })

                        # 添加分隔线
                        merged_doc.add_paragraph('─' * 50)

                    except Exception as e:
                        error_msg = f'处理错误: {str(e)}'
                        print(f"处理文件 {file.original_name} 时出错: {error_msg}")
                        skipped_files.append({
                            'name': file.original_name,
                            'reason': error_msg
                        })

            # 在每个阶段后添加空白页，标题为阶段名
            print(f"为阶段 {stage.name} 添加分隔页")
            merged_doc.add_page_break()
            blank_page = merged_doc.add_heading(f'{stage.name}', level=1)
            blank_page.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
            # 添加页面分隔符，开始新的阶段
            merged_doc.add_page_break()

        # 添加统计信息
        print("添加统计信息")
        merged_doc.add_heading('文档合并统计', level=1)
        merged_doc.add_paragraph(f'成功处理文件数: {len(processed_files)}')
        if processed_files:
            merged_doc.add_paragraph('成功处理的文件:')
            for filename in processed_files:
                merged_doc.add_paragraph(filename, style='List Bullet')

        if skipped_files:
            merged_doc.add_paragraph(f'跳过文件数: {len(skipped_files)}')
            merged_doc.add_paragraph('跳过的文件:')
            for file_info in skipped_files:
                p = merged_doc.add_paragraph('', style='List Bullet')
                p.add_run(f"{file_info['name']} - 原因: {file_info['reason']}")

        # 保存合并后的文档
        output_filename = f"{project.name}_merged.docx"
        output_path = os.path.join(temp_dir, output_filename)
        print(f"保存合并文档到: {output_path}")

        try:
            merged_doc.save(output_path)
        except Exception as e:
            error_msg = f"保存合并文档失败: {str(e)}"
            print(error_msg)
            return None, error_msg

        # 检查文件是否成功创建
        if not os.path.exists(output_path):
            error_msg = "生成的文档不存在"
            print(error_msg)
            return None, error_msg

        if os.path.getsize(output_path) == 0:
            error_msg = "生成的文档为空"
            print(error_msg)
            return None, error_msg

        print(f"文档已成功保存，大小: {os.path.getsize(output_path)} 字节")

        # 创建临时文件返回
        try:
            print("创建返回用的临时文件")
            temp_output = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
            temp_output.close()  # 关闭文件句柄

            print(f"复制文件 {output_path} 到 {temp_output.name}")
            shutil.copy2(output_path, temp_output.name)

            if not os.path.exists(temp_output.name):
                error_msg = "复制到临时文件失败: 目标文件不存在"
                print(error_msg)
                return None, error_msg

            if os.path.getsize(temp_output.name) == 0:
                error_msg = "复制到临时文件失败: 目标文件为空"
                print(error_msg)
                return None, error_msg

            print(f"文件成功复制到临时位置，大小: {os.path.getsize(temp_output.name)} 字节")
            return temp_output.name, None
        except Exception as e:
            error_msg = f"创建临时文件失败: {str(e)}"
            print(error_msg)
            return None, error_msg

    except Exception as e:
        error_msg = f"合并过程中发生错误: {str(e)}"
        print(error_msg)
        import traceback
        traceback.print_exc()
        return None, error_msg
    finally:
        if temp_dir and os.path.exists(temp_dir):
            try:
                print(f"清理临时目录: {temp_dir}")
                shutil.rmtree(temp_dir)
            except Exception as e:
                print(f"清理临时文件失败: {str(e)}")