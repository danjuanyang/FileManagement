import base64
import re
import subprocess
import tempfile
import urllib
import os
import uuid

import io
from flask import Blueprint, request, jsonify, current_app, send_from_directory
from datetime import datetime
import os

from pptx.enum.shapes import MSO_SHAPE_TYPE
from sqlalchemy.exc import SQLAlchemyError
from werkzeug.utils import secure_filename
from models import db, Training, Comment, Reply, User
from routes.employees import token_required
from routes.filemanagement import python_dir
from utils.activity_tracking import track_activity
import urllib.parse

import tempfile
from pptx import Presentation
import comtypes.client
from PIL import Image, ImageDraw, ImageFont
from pptx.enum.dml import MSO_THEME_COLOR_INDEX, MSO_FILL_TYPE
from pptx.enum.dml import MSO_FILL

# 创建蓝图
training_bp = Blueprint('training', __name__)

# 文件上传配置
ALLOWED_EXTENSIONS = {'ppt', 'pptx'}
UPLOAD_FOLDER = os.path.join(python_dir, 'uploads')  # 基础上传目录


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def sanitize_filename(filename):
    """
    安全地处理文件名，保留中文字符
    移除或替换不安全的字符，但保留中文和基本标点
    """
    # 替换Windows文件系统中的非法字符
    illegal_chars = r'[<>:"/\\|?*\x00-\x1f]'
    # 将非法字符替换为下划线
    safe_name = re.sub(illegal_chars, '_', filename)
    # 去除首尾空格和点
    safe_name = safe_name.strip('. ')
    # 如果文件名为空，返回默认名称
    return safe_name or 'unnamed'


def get_safe_path_component(component):
    """
    安全地处理路径组件名称，保留中文字符
    """
    if not component:
        return 'unnamed'
    return sanitize_filename(component)


def generate_unique_filename(directory, original_filename):
    """生成唯一的文件名，处理文件名冲突"""
    base_name, extension = os.path.splitext(original_filename)
    safe_base_name = get_safe_path_component(base_name)
    counter = 1
    new_filename = f"{safe_base_name}{extension}"

    while os.path.exists(os.path.join(directory, new_filename)):
        new_filename = f"{safe_base_name}({counter}){extension}"
        counter += 1

    return new_filename


def ensure_upload_folder():
    """确保上传目录存在"""
    folder_path = os.path.join(current_app.static_folder, UPLOAD_FOLDER)
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)
    return folder_path


# 分配培训任务
@training_bp.route('/assign', methods=['POST'])
@token_required
@track_activity
def assign_training(current_user):
    try:
        data = request.get_json()
        if not data:
            return jsonify({
                'code': 400,
                'message': '无效的请求数据'
            }), 400

        # 获取并验证必要字段
        trainer_id = data.get('trainer_id')
        training_month = data.get('training_month')
        title = data.get('title')
        description = data.get('description')

        # 检查必要字段
        if not all([trainer_id, training_month, title]):
            return jsonify({
                'code': 400,
                'message': '缺少必要字段'
            }), 400

        # 确保trainer_id是整数类型
        try:
            trainer_id = int(trainer_id)
        except (TypeError, ValueError):
            return jsonify({
                'code': 400,
                'message': '无效的培训人ID'
            }), 400

        # 验证培训人是否存在
        trainer = User.query.get(trainer_id)
        if not trainer:
            return jsonify({
                'code': 404,
                'message': '培训人不存在'
            }), 404

        # 验证月份格式 (YYYY-MM)
        try:
            datetime.strptime(training_month, '%Y-%m')
        except ValueError:
            return jsonify({
                'code': 400,
                'message': '无效的月份格式，应为YYYY-MM'
            }), 400

        # 检查是否已存在该月份的培训任务
        existing_training = Training.query.filter_by(
            trainer_id=trainer_id,
            training_month=training_month
        ).first()

        if existing_training:
            return jsonify({
                'code': 400,
                'message': '该员工在此月份已有培训任务'
            }), 400

        # 创建新的培训任务
        new_training = Training(
            trainer_id=trainer_id,  # 确保这里使用正确的trainer_id
            training_month=training_month,
            title=title,
            description=description,
            status='pending',  # 初始状态为待上传
            create_time=datetime.now()
        )

        db.session.add(new_training)
        db.session.commit()

        # 返回成功响应
        return jsonify({
            'code': 200,
            'message': '培训任务分配成功',
            'data': {
                'id': new_training.id,
                'trainer_id': new_training.trainer_id,
                'trainer_name': trainer.username,  # 添加培训人姓名
                'training_month': new_training.training_month,
                'title': new_training.title,
                'status': new_training.status
            },
        })

    except Exception as e:
        db.session.rollback()
        print(f"Error in assign_training: {str(e)}")  # 添加错误日志
        return jsonify({
            'code': 500,
            'message': f'服务器错误: {str(e)}'
        }), 500


# 上传培训材料
@training_bp.route('/upload/<int:training_id>', methods=['POST'])
@token_required
@track_activity
def upload_training_material(current_user, training_id):
    try:
        # 获取培训记录
        training = Training.query.get_or_404(training_id)

        # 验证当前用户是否为指定培训者
        if current_user.id != training.trainer_id:
            return jsonify({
                'code': 403,
                'message': '您不是该培训的指定培训者'
            }), 403

        # 验证培训状态是否为待上传
        if training.status != 'pending':
            return jsonify({
                'code': 400,
                'message': '该培训材料已上传完成，无法重复上传'
            }), 400

        # 检查是否有文件被上传
        if 'file' not in request.files:
            return jsonify({
                'code': 400,
                'message': '未找到上传文件'
            }), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({
                'code': 400,
                'message': '未选择文件'
            }), 400

        # 验证文件类型
        if not allowed_file(file.filename):
            return jsonify({
                'code': 400,
                'message': '不支持的文件类型，仅支持 PPT/PPTX 格式'
            }), 400

        # 确保上传目录存在
        upload_folder = os.path.join(current_app.root_path, 'uploads', 'training')
        os.makedirs(upload_folder, exist_ok=True)

        # 处理文件名，保留中文字符并确保唯一性
        original_filename = file.filename
        safe_filename = generate_unique_filename(upload_folder, original_filename)

        # 保存文件
        file_path = os.path.join(upload_folder, safe_filename)
        file.save(file_path)

        # 更新培训记录，存储相对路径
        training.material_path = os.path.join('training', safe_filename)
        training.status = 'completed'
        training.upload_time = datetime.now()

        db.session.commit()

        return jsonify({
            'code': 200,
            'message': '培训材料上传成功',
            'data': {
                'original_name': original_filename,
                'saved_name': safe_filename,
                'upload_time': training.upload_time.strftime('%Y-%m-%d %H:%M:%S'),
                'status': training.status
            }
        })

    except Exception as e:
        db.session.rollback()
        print(f"upload_training_material 错误： {str(e)}")
        return jsonify({
            'code': 500,
            'message': f'服务器错误: {str(e)}'
        }), 500


# 获取培训列表
@training_bp.route('/list', methods=['GET'])
@token_required
@track_activity
def get_training_list(current_user):  # 添加 current_user 参数
    try:
        page = request.args.get('page', 1, type=int)
        per_page = request.args.get('per_page', 10, type=int)
        month = request.args.get('month')  # 可选过滤条件

        query = Training.query

        if month:
            query = query.filter_by(training_month=month)

        trainings = query.order_by(Training.training_month.desc()) \
            .paginate(page=page, per_page=per_page, error_out=False)

        return jsonify({
            'code': 200,
            'data': {
                'items': [{
                    'id': t.id,
                    'trainer_id': t.trainer_id,
                    'trainer_name': t.trainer.username,  # 假设User模型中有name字段
                    'training_month': t.training_month,
                    'title': t.title,
                    'description': t.description,
                    'status': t.status,
                    'material_path': t.material_path,
                    'upload_time': t.upload_time.strftime('%Y-%m-%d %H:%M:%S') if t.upload_time else None
                } for t in trainings.items],
                'total': trainings.total,
                'pages': trainings.pages,
                'current_page': trainings.page
            }
        })

    except Exception as e:
        print(f"Error in get_training_list: {str(e)}")  # 添加日志输出
        return jsonify({
            'code': 500,
            'message': f'服务器错误: {str(e)}'
        }), 500


# 获取培训详情
@training_bp.route('/<int:training_id>', methods=['GET'])
@token_required
@track_activity
def get_training_detail(current_user, training_id):
    try:
        training = Training.query.get_or_404(training_id)

        # 获取评论及其回复
        comments = Comment.query.filter_by(training_id=training_id) \
            .order_by(Comment.create_time.desc()).all()

        return jsonify({
            'code': 200,
            'data': {
                'id': training.id,
                'trainer_id': training.trainer_id,
                'trainer_name': training.trainer.username,
                'training_month': training.training_month,
                'title': training.title,
                'description': training.description,
                'status': training.status,
                'material_path': training.material_path,
                'is_trainer': current_user.id == training.trainer_id,
                'upload_time': training.upload_time.strftime('%Y-%m-%d %H:%M:%S') if training.upload_time else None,
                'comments': [{
                    'id': comment.id,
                    'user_id': comment.user_id,
                    'user_name': comment.user.name,
                    'content': comment.content,
                    'create_time': comment.create_time.strftime('%Y-%m-%d %H:%M:%S'),
                    'replies': [{
                        'id': reply.id,
                        'user_id': reply.user_id,
                        'user_name': reply.user.name,
                        'content': reply.content,
                        'create_time': reply.create_time.strftime('%Y-%m-%d %H:%M:%S')
                    } for reply in comment.replies]
                } for comment in comments]
            }
        })

    except Exception as e:
        return jsonify({
            'code': 500,
            'message': f'服务器错误: {str(e)}'
        }), 500


# 添加评论
@training_bp.route('/<int:training_id>/comment', methods=['POST'])
@token_required
@track_activity
def add_comment(training_id):
    try:
        data = request.get_json()
        content = data.get('content')

        if not content:
            return jsonify({
                'code': 400,
                'message': '评论内容不能为空'
            }), 400

        comment = Comment(
            training_id=training_id,
            user_id=request.user.id,
            content=content
        )

        db.session.add(comment)
        db.session.commit()

        return jsonify({
            'code': 200,
            'message': '评论成功',
            'data': {
                'id': comment.id,
                'user_id': comment.user_id,
                'user_name': comment.user.name,
                'content': comment.content,
                'create_time': comment.create_time.strftime('%Y-%m-%d %H:%M:%S')
            }
        })

    except Exception as e:
        db.session.rollback()
        return jsonify({
            'code': 500,
            'message': f'服务器错误: {str(e)}'
        }), 500


# 添加回复
@training_bp.route('/comment/<int:comment_id>/reply', methods=['POST'])
@token_required
@track_activity
def add_reply(comment_id):
    try:
        data = request.get_json()
        content = data.get('content')

        if not content:
            return jsonify({
                'code': 400,
                'message': '回复内容不能为空'
            }), 400

        reply = Reply(
            comment_id=comment_id,
            user_id=request.user.id,
            content=content
        )

        db.session.add(reply)
        db.session.commit()

        return jsonify({
            'code': 200,
            'message': '回复成功',
            'data': {
                'id': reply.id,
                'user_id': reply.user_id,
                'user_name': reply.user.name,
                'content': reply.content,
                'create_time': reply.create_time.strftime('%Y-%m-%d %H:%M:%S')
            }
        })

    except Exception as e:
        db.session.rollback()
        return jsonify({
            'code': 500,
            'message': f'服务器错误: {str(e)}'
        }), 500

# 预览PPT
@training_bp.route('/preview/<path:filename>')
@token_required
def preview_training_file(current_user, filename):
    try:
        # 解码文件名
        filename = urllib.parse.unquote(filename)

        # 构建完整的文件路径
        file_path = os.path.join(current_app.root_path, 'uploads', 'training', os.path.basename(filename))

        if not os.path.exists(file_path):
            return jsonify({
                'code': 404,
                'message': '文件不存在'
            }), 404

        # 使用纯Python处理器处理PPT
        handler = PurePPTHandler()
        result = handler.process_ppt(file_path)

        if not result['success']:
            return jsonify({
                'code': 500,
                'message': f'处理PPT失败: {result["error"]}'
            }), 500

        return jsonify({
            'code': 200,
            'data': {
                'slides': result['slides'],
                'total_slides': result['total_slides']
            }
        })

    except Exception as e:
        return jsonify({
            'code': 500,
            'message': f'预览生成失败: {str(e)}'
        }), 500


class PurePPTHandler:
    @staticmethod
    def get_system_font():
        """获取系统中文字体"""
        # 常见的中文字体路径
        font_paths = [
            # Windows 字体
            "C:/Windows/Fonts/simhei.ttf",  # 黑体
            "C:/Windows/Fonts/simsun.ttc",  # 宋体
            "C:/Windows/Fonts/msyh.ttc",  # 微软雅黑
            # Linux 字体
            "/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf",
            "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
            # macOS 字体
            "/System/Library/Fonts/PingFang.ttc",
            # 回退选项
            "",  # 空字符串会触发默认字体
        ]

        for font_path in font_paths:
            try:
                if os.path.exists(font_path):
                    return ImageFont.truetype(font_path, 24)
            except Exception:
                continue

        return ImageFont.load_default()

    @staticmethod
    def extract_ppt_content(ppt_path):
        """
        提取PPT内容，将每页转换为简单的文本和形状表示
        """
        print(f"Processing PPT file: {ppt_path}")
        prs = Presentation(ppt_path)
        slides_data = []

        for i, slide in enumerate(prs.slides):
            print(f"Processing slide {i + 1}")
            slide_data = {
                'shapes': [],
                'texts': [],
                'background_color': '#FFFFFF'  # 默认白色背景
            }

            # 提取背景（安全处理）
            try:
                background = slide.background
                if hasattr(background, 'fill') and background.fill:
                    fill_type = getattr(background.fill, 'type', None)
                    if fill_type is not None and fill_type != MSO_FILL_TYPE.NONE:
                        if hasattr(background.fill, 'fore_color') and background.fill.fore_color:
                            rgb = background.fill.fore_color.rgb
                            if rgb:
                                slide_data['background_color'] = f'#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}'
            except Exception as e:
                print(f"背景提取警告（非严重）： {e}")
                # 继续使用默认白色背景

            # 提取形状和文本
            for shape in slide.shapes:
                try:
                    # 处理文本
                    if hasattr(shape, "text") and shape.text.strip():
                        # 获取文本框的填充颜色
                        text_color = '#000000'  # 默认黑色
                        try:
                            if shape.fill.type != MSO_FILL.NONE:
                                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                                    rgb = shape.fill.fore_color.rgb
                                    if rgb:
                                        text_color = f'#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}'
                        except Exception:
                            pass

                        slide_data['texts'].append({
                            'text': shape.text,
                            'left': shape.left / 9144000 * 1024,  # 转换EMU到像素
                            'top': shape.top / 6858000 * 768,
                            'width': shape.width / 9144000 * 1024,
                            'height': shape.height / 6858000 * 768,
                            'color': text_color
                        })

                    # 处理图片
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            print(f"Found picture shape at position ({shape.left}, {shape.top})")
                            # 直接获取图片数据
                            image = shape.image
                            if not hasattr(image, 'blob'):
                                print("Image has no blob attribute")
                                continue

                            image_bytes = image.blob
                            if not image_bytes:
                                print("Image blob is empty")
                                continue

                            print(f"Image format: {getattr(image, 'content_type', 'unknown')}")
                            print(f"Image size: {len(image_bytes)} bytes")

                            # 使用PIL验证图片数据
                            try:
                                test_image = Image.open(io.BytesIO(image_bytes))
                                print(f"Image mode: {test_image.mode}, size: {test_image.size}")
                            except Exception as img_err:
                                print(f"Invalid image data: {img_err}")
                                continue

                            # 转换为base64
                            img_base64 = base64.b64encode(image_bytes).decode()

                            # 计算位置（添加偏移修正）
                            left = max(0, shape.left / 9144000 * 1024)
                            top = max(0, shape.top / 6858000 * 768)
                            width = min(1024, shape.width / 9144000 * 1024)
                            height = min(768, shape.height / 6858000 * 768)

                            slide_data['shapes'].append({
                                'type': 'image',
                                'data': img_base64,
                                'left': left,
                                'top': top,
                                'width': width,
                                'height': height,
                                'format': getattr(image, 'content_type', 'png')
                            })
                            print(f"Successfully added image: {left:.0f}x{top:.0f} @ {width:.0f}x{height:.0f}")
                        except Exception as e:
                            print(f"Image processing error: {str(e)}")
                            import traceback
                            print(traceback.format_exc())
                except Exception as e:
                    print(f"Shape processing error: {e}")

            slides_data.append(slide_data)

        return slides_data

    @staticmethod
    def render_slide_preview(slide_data, width=1024, height=768):
        """
        将幻灯片数据渲染为简单的预览图
        """
        try:
            # 创建带背景色的图片
            background_color = slide_data.get('background_color', '#FFFFFF')
            image = Image.new('RGB', (width, height), background_color)
            draw = ImageDraw.Draw(image)

            # 获取系统中文字体
            font = PurePPTHandler.get_system_font()

            # 渲染形状
            print(f"Processing {len(slide_data['shapes'])} shapes")
            for i, shape in enumerate(slide_data['shapes']):
                if shape['type'] == 'image':
                    try:
                        print(f"Rendering image {i + 1}")
                        # 解码base64
                        img_data = base64.b64decode(shape['data'])
                        print(f"Decoded image size: {len(img_data)} bytes")

                        # 打开图片
                        shape_image = Image.open(io.BytesIO(img_data))
                        print(f"Original image mode: {shape_image.mode}, size: {shape_image.size}")

                        # 确保图片大小合理
                        target_width = max(1, min(int(shape['width']), width))
                        target_height = max(1, min(int(shape['height']), height))
                        print(f"Resizing to: {target_width}x{target_height}")

                        # 调整图片大小
                        shape_image = shape_image.resize(
                            (target_width, target_height),
                            Image.Resampling.LANCZOS
                        )

                        # 确保图片模式正确
                        if shape_image.mode in ['RGBA', 'LA']:
                            # 创建白色背景
                            background = Image.new('RGBA', shape_image.size, 'white')
                            shape_image = Image.alpha_composite(background.convert('RGBA'),
                                                                shape_image.convert('RGBA'))

                        # 转换为RGB模式
                        if shape_image.mode != 'RGB':
                            shape_image = shape_image.convert('RGB')

                        # 计算粘贴位置
                        paste_x = max(0, min(int(shape['left']), width - target_width))
                        paste_y = max(0, min(int(shape['top']), height - target_height))
                        print(f"Pasting at position: ({paste_x}, {paste_y})")

                        # 粘贴图片
                        try:
                            image.paste(
                                shape_image,
                                (paste_x, paste_y)
                            )
                            print("Successfully pasted image")
                        except Exception as paste_err:
                            print(f"Error pasting image: {paste_err}")

                    except Exception as e:
                        print(f"Error rendering image shape: {str(e)}")
                        import traceback
                        print(traceback.format_exc())

            # 渲染文本
            for text_item in slide_data['texts']:
                try:
                    x = int(text_item['left'])
                    y = int(text_item['top'])

                    # 文本换行处理
                    max_width = int(text_item['width'])
                    text = text_item['text']
                    color = text_item.get('color', '#000000')

                    # 简单的文本换行
                    words = text.split()
                    lines = []
                    current_line = []

                    for word in words:
                        current_line.append(word)
                        test_line = ' '.join(current_line)
                        bbox = draw.textbbox((0, 0), test_line, font=font)
                        if bbox[2] - bbox[0] > max_width:
                            current_line.pop()
                            lines.append(' '.join(current_line))
                            current_line = [word]

                    lines.append(' '.join(current_line))

                    # 绘制每行文本
                    line_height = font.size + 4
                    for i, line in enumerate(lines):
                        draw.text(
                            (x, y + i * line_height),
                            line,
                            fill=color,
                            font=font
                        )

                except Exception as e:
                    print(f"Error rendering text: {e}")

            # 保存为PNG
            buffer = io.BytesIO()
            image.save(buffer, format='PNG', optimize=True)
            buffer.seek(0)
            return base64.b64encode(buffer.getvalue()).decode()

        except Exception as e:
            print(f"Error in render_slide_preview: {e}")
            return None

    @staticmethod
    def process_ppt(ppt_path):
        """
        处理PPT文件并返回所有幻灯片的预览
        """
        try:
            print(f"Starting to process PPT: {ppt_path}")
            slides_data = PurePPTHandler.extract_ppt_content(ppt_path)

            previews = []
            for i, slide_data in enumerate(slides_data):
                print(f"Rendering slide {i + 1}")
                preview = PurePPTHandler.render_slide_preview(slide_data)
                if preview:
                    previews.append({
                        'image': preview,
                        'texts': slide_data['texts'],
                        'shapes': [s for s in slide_data['shapes'] if s['type'] == 'image']
                    })
                else:
                    print(f"Failed to render slide {i + 1}")

            return {
                'success': True,
                'slides': previews,
                'total_slides': len(previews)
            }

        except Exception as e:
            print(f"Error in process_ppt: {e}")
            return {
                'success': False,
                'error': str(e)
            }


# 下载
@training_bp.route('/download/<path:filename>')
@token_required
def download_training_file(current_user, filename):
    try:
        # 解码文件名
        filename = urllib.parse.unquote(filename)

        # 构建完整的文件路径
        file_path = os.path.join(current_app.root_path, 'uploads', 'training', os.path.basename(filename))

        if not os.path.exists(file_path):
            return jsonify({
                'code': 404,
                'message': '文件不存在'
            }), 404

        # 获取文件的原始名称
        directory = os.path.dirname(file_path)
        return send_from_directory(
            directory,
            os.path.basename(file_path),
            as_attachment=True
        )

    except Exception as e:
        return jsonify({
            'code': 500,
            'message': f'下载失败: {str(e)}'
        }), 500